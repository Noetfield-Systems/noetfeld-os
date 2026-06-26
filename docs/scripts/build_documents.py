#!/usr/bin/env python3
"""
Generate Noetfield internal / external PDFs and pitch-deck slide files.

Usage:
    pip install fpdf2 python-pptx
    python docs/scripts/build_documents.py
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import date
from pathlib import Path

from fpdf import FPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "output"


def _ascii(text: str) -> str:
    """Normalize text for Helvetica PDF output."""
    return (
        text.replace("\u2014", " - ")
        .replace("\u2013", "-")
        .replace("\u2192", "->")
        .replace("\u2714", "[x]")
        .replace("\u2022", "-")
    )

NOETFIELD_BLUE = (30, 64, 120)
NOETFIELD_DARK = (18, 32, 58)
ACCENT = (56, 120, 210)
MUTED = (90, 100, 115)


@dataclass(frozen=True)
class DocMeta:
    filename: str
    title: str
    subtitle: str
    audience: str  # INTERNAL | EXTERNAL
    classification: str


class NoetfieldPDF(FPDF):
    def __init__(self, meta: DocMeta) -> None:
        super().__init__()
        self.meta = meta
        self.set_auto_page_break(auto=True, margin=18)
        self.set_margins(18, 18, 18)

    def header(self) -> None:
        self.set_fill_color(*NOETFIELD_BLUE)
        self.rect(0, 0, 210, 22, style="F")
        self.set_xy(18, 7)
        self.set_font("Helvetica", "B", 11)
        self.set_text_color(255, 255, 255)
        self.cell(0, 8, "NOETFIELD SYSTEMS INC.", new_x="LMARGIN", new_y="NEXT")
        self.set_xy(18, 28)
        self.set_text_color(*NOETFIELD_DARK)

    def footer(self) -> None:
        self.set_y(-14)
        self.set_font("Helvetica", "", 8)
        self.set_text_color(*MUTED)
        label = f"{self.meta.audience} | {self.meta.classification}"
        self.cell(0, 8, f"{label}  |  Page {self.page_no()}/{{nb}}  |  noetfield.com", align="C")

    def cover(self) -> None:
        self.add_page()
        self.set_y(45)
        self.set_font("Helvetica", "B", 22)
        self.set_text_color(*NOETFIELD_DARK)
        self.multi_cell(self.epw, 11, _ascii(self.meta.title))
        self.ln(4)
        self.set_font("Helvetica", "", 13)
        self.set_text_color(*MUTED)
        self.multi_cell(self.epw, 8, _ascii(self.meta.subtitle))
        self.ln(8)
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(*ACCENT)
        self.cell(self.epw, 8, f"Audience: {self.meta.audience}", new_x="LMARGIN", new_y="NEXT")
        self.cell(self.epw, 8, f"Classification: {self.meta.classification}", new_x="LMARGIN", new_y="NEXT")
        self.set_font("Helvetica", "", 10)
        self.set_text_color(*MUTED)
        self.cell(self.epw, 8, f"Generated: {date.today().isoformat()}", new_x="LMARGIN", new_y="NEXT")
        self.ln(10)
        self.set_draw_color(*NOETFIELD_BLUE)
        y = self.get_y()
        self.line(self.l_margin, y, self.w - self.r_margin, y)

    def section(self, title: str) -> None:
        self.set_x(self.l_margin)
        self.ln(4)
        self.set_font("Helvetica", "B", 13)
        self.set_text_color(*NOETFIELD_BLUE)
        self.multi_cell(self.epw, 8, _ascii(title))
        self.ln(2)

    def body(self, text: str) -> None:
        self.set_x(self.l_margin)
        self.set_font("Helvetica", "", 10.5)
        self.set_text_color(35, 40, 48)
        self.multi_cell(self.epw, 6, _ascii(text))
        self.ln(2)

    def bullets(self, items: list[str]) -> None:
        self.set_x(self.l_margin)
        self.set_font("Helvetica", "", 10.5)
        self.set_text_color(35, 40, 48)
        for item in items:
            self.multi_cell(self.epw, 6, _ascii(f"  -  {item}"))
        self.ln(2)

    def table(self, headers: list[str], rows: list[list[str]], col_widths: list[int] | None = None) -> None:
        self.set_x(self.l_margin)
        if col_widths is None:
            width = int(self.epw // len(headers))
            col_widths = [width] * len(headers)
        self.set_font("Helvetica", "B", 9)
        self.set_fill_color(230, 236, 246)
        for i, h in enumerate(headers):
            self.cell(col_widths[i], 8, _ascii(h), border=1, fill=True)
        self.ln()
        self.set_font("Helvetica", "", 9)
        for row in rows:
            for i, cell in enumerate(row):
                self.cell(col_widths[i], 8, _ascii(cell), border=1)
            self.ln()
        self.ln(3)


def save_pdf(meta: DocMeta, builder) -> Path:
    pdf = NoetfieldPDF(meta)
    pdf.alias_nb_pages()
    builder(pdf)
    out_dir = OUT / meta.audience.lower()
    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / meta.filename
    pdf.output(str(path))
    return path


# ---------------------------------------------------------------------------
# PDF builders
# ---------------------------------------------------------------------------

def build_grant_narrative_external(pdf: NoetfieldPDF) -> None:
    pdf.cover()
    pdf.section("Executive Positioning")
    pdf.body(
        "Noetfield is a pre-execution governance infrastructure that enables regulated "
        "organizations to simulate, evaluate, and audit AI-driven operational decisions "
        "before deployment."
    )
    pdf.section("Problem Statement")
    pdf.body(
        "Canadian regulated organizations are rapidly adopting AI systems without "
        "pre-execution validation, audit-grade decision traceability, policy enforcement "
        "before system actions, or regulatory-ready evidence generation."
    )
    pdf.bullets([
        "Compliance risk (AI Act and emerging AI governance rules)",
        "Operational blind spots in automated decision chains",
        "Lack of explainability for board-level oversight",
    ])
    pdf.section("Solution: Governance Execution Layer (GEL)")
    pdf.bullets([
        "Intent evaluation before execution (ALLOW / REJECT / REVIEW)",
        "Immutable audit ledger with tamper-evident decision records",
        "Policy enforcement engine (Golden Edge v3 governance rules)",
        "Simulation console for regulated decision testing",
    ])
    pdf.section("Innovation - Why This Warrants Public Support")
    pdf.body(
        "This project introduces structural separation of execution vs governance. "
        "Noetfield NEVER executes financial or operational actions. It ONLY evaluates, "
        "simulates, and logs decision outcomes while producing audit-ready compliance evidence."
    )
    pdf.section("Strategic Value to Canada")
    pdf.bullets([
        "AI governance infrastructure aligned with emerging regulatory frameworks",
        "Reduces enterprise AI risk exposure in banking and public sector",
        "Enables safe adoption of AI without expanding custody or settlement scope",
        "Strengthens Canadian AI compliance tooling ecosystem",
    ])


def build_grant_narrative_internal(pdf: NoetfieldPDF) -> None:
    build_grant_narrative_external(pdf)
    pdf.add_page()
    pdf.section("Internal Appendix - Build Status vs Narrative")
    pdf.body(
        "Use external narrative for IRAP and partner conversations. This appendix "
        "maps claims to current engineering reality in NOETFELD OS."
    )
    pdf.table(
        ["Capability", "External claim", "Current build", "Target (pilot)"],
        [
            ["Decision engine", "ALLOW/REJECT/REVIEW", "Working (FastAPI)", "Production hardened"],
            ["Policy layer", "Golden Edge v3", "JSON policy loader", "Versioned tenant policies"],
            ["Audit ledger", "Postgres-backed", "SQLite audit_log", "Postgres migration"],
            ["Console", "Simulation UI", "API + /portal/audits", "Full simulation console"],
            ["Non-custodial", "No execution", "Decision-only API", "Contract + docs aligned"],
        ],
        [42, 38, 44, 50],
    )
    pdf.section("IRAP Messaging Guardrails")
    pdf.bullets([
        "Lead with governance infrastructure, not lending or payments.",
        "Emphasize simulation and audit — not autonomous execution.",
        "Position Golden Edge v3 as policy framework name; document rule mapping.",
        "Pilot scope: read-only integration, 6 weeks, BC organization.",
    ])


def build_loi_external(pdf: NoetfieldPDF) -> None:
    pdf.cover()
    pdf.section("Letter of Intent")
    pdf.body("To: Noetfield Systems Inc.")
    pdf.body("From: [Organization Name]")
    pdf.body("Date: [Insert Date]")
    pdf.body("Subject: Intent to Participate in Noetfield Governance Execution Pilot")
    pdf.ln(2)
    pdf.body(
        "We hereby express interest in participating in the Noetfield Governance "
        "Execution Pilot Program."
    )
    pdf.section("We Understand That")
    pdf.bullets([
        "The system operates as a pre-execution governance and audit layer",
        "No financial custody or transaction processing is performed",
        "The pilot is strictly for evaluation, simulation, and compliance readiness",
    ])
    pdf.section("Scope of Engagement")
    pdf.bullets([
        "6-week governance assessment",
        "AI decision simulation testing",
        "Policy enforcement mapping",
        "Audit trail review",
    ])
    pdf.section("Organizational Interest")
    pdf.bullets([
        "[ ] AI risk reduction",
        "[ ] Compliance readiness",
        "[ ] Governance modernization",
        "[ ] Board-level audit visibility",
    ])
    pdf.section("Non-Binding Agreement")
    pdf.body(
        "This letter does not constitute a financial commitment but indicates "
        "organizational intent to explore collaboration."
    )
    pdf.ln(8)
    pdf.body("Authorized Signatory: ___________________________")
    pdf.body("Title: ___________________________")
    pdf.body("Organization: ___________________________")


def build_loi_internal(pdf: NoetfieldPDF) -> None:
    build_loi_external(pdf)
    pdf.add_page()
    pdf.section("Internal Sales Notes (Do Not Share With Prospect)")
    pdf.bullets([
        "Use after first discovery call; not a substitute for MSA or SOW.",
        "Target: BC regulated org, bank-adjacent, or enterprise AI team.",
        "Follow-up within 48h: send Trust Brief one-pager + pilot timeline.",
        "Convert LOI to paid pilot ($10K entry) or IRAP co-funded engagement ($120K).",
        "Procurement path: procurement@noetfield.com for invoice/PO.",
    ])


def build_budget_external(pdf: NoetfieldPDF) -> None:
    pdf.cover()
    pdf.section("Project Overview")
    pdf.body("Total Project Cost: CAD $120,000 | Duration: 6-week pilot program")
    pdf.section("Cost Structure")
    pdf.table(
        ["Category", "Description", "Cost (CAD)"],
        [
            ["Core Engineering", "Governance engine + API refinement", "$35,000"],
            ["Compliance Modeling", "Policy mapping + regulatory alignment", "$20,000"],
            ["Data Infrastructure", "Audit + event ledger system", "$15,000"],
            ["Pilot Deployment", "Client integration + sandbox setup", "$20,000"],
            ["Evaluation & Reporting", "Trust Brief + audit report generation", "$20,000"],
            ["Overhead & Research", "Architecture + system refinement", "$10,000"],
        ],
        [48, 86, 40],
    )
    pdf.section("Funding Request")
    pdf.body(
        "Requested support from IRAP / BC Innovation: CAD $60,000 "
        "(50% co-funding model). Remainder from pilot partner fees and Noetfield contribution."
    )


def build_budget_internal(pdf: NoetfieldPDF) -> None:
    build_budget_external(pdf)
    pdf.section("Internal — Commercial Ladder")
    pdf.table(
        ["Tier", "Price", "Scope", "Purpose"],
        [
            ["Entry pilot", "$10,000", "2-week sandbox + policy review", "Land first logo"],
            ["Standard pilot", "$50,000", "6-week GEL pilot + Trust Brief", "Core revenue target"],
            ["IRAP program", "$120,000", "Full co-funded build + enterprise pilot", "Grant + validation"],
        ],
        [36, 28, 62, 48],
    )
    pdf.section("Internal — Delivery Assumptions")
    pdf.bullets([
        "Engineering headcount: 1 FTE + founder oversight for 6 weeks.",
        "Postgres migration budgeted in Data Infrastructure line.",
        "Trust Brief deliverable reuses existing playbook templates.",
        "50% IRAP match requires documented eligible R&D activities.",
    ])


def build_strategic_external(pdf: NoetfieldPDF) -> None:
    pdf.cover()
    pdf.section("Core Positioning Statement")
    pdf.body(
        "Noetfield is not an AI platform. It is a governance execution and audit "
        "infrastructure for regulated AI systems."
    )
    pdf.section("What We Are")
    pdf.bullets([
        "Pre-execution governance layer for AI-driven operational decisions",
        "Policy validation and corridor enforcement before external systems act",
        "Audit-grade decision traces suitable for board and regulatory review",
        "Non-custodial by design — no funds, no settlement, no payment routing",
    ])
    pdf.section("What We Are Not")
    pdf.bullets([
        "Not a general-purpose AI model provider",
        "Not a lending platform or payment processor",
        "Not an autonomous execution engine",
    ])
    pdf.section("Ideal Customer Profile")
    pdf.bullets([
        "Canadian regulated enterprise (banking, insurance, public sector)",
        "Teams deploying AI for operational decisions requiring explainability",
        "Organizations preparing for AI governance and audit scrutiny",
    ])


def build_strategic_internal(pdf: NoetfieldPDF) -> None:
    build_strategic_external(pdf)
    pdf.add_page()
    pdf.section("Internal — Competitive Moat")
    pdf.bullets([
        "Execution/governance separation is the category — not feature parity with LLM vendors.",
        "Trust Ledger + GEL bundle creates recurring governance revenue.",
        "IRAP narrative opens non-dilutive capital while building enterprise proof.",
        "BC base + Canadian regulatory framing = differentiated vs US-only tooling.",
    ])
    pdf.section("Internal — Messaging Do / Don't")
    pdf.table(
        ["Do", "Don't"],
        [
            ["Governance execution layer", "AI platform"],
            ["Simulate and audit decisions", "Approve loans or move money"],
            ["Board-ready evidence", "Replace compliance officers"],
            ["Read-only pilot integration", "Production execution authority"],
        ],
        [87, 87],
    )


def build_reality_check_internal(pdf: NoetfieldPDF) -> None:
    pdf.cover()
    pdf.section("What We Have Today (Engineering Reality)")
    pdf.table(
        ["Asset", "Status", "Location / Notes"],
        [
            ["Governance engine", "Working", "decision_engine.py + policy_loader.py"],
            ["Decision API", "Working", "POST /v1/decision"],
            ["Audit system", "Working (SQLite)", "database.py — migrate to Postgres"],
            ["Portal / audit read", "Working", "GET /portal/audits"],
            ["Policy files", "Working", "base_policy.json, corridor_policy.json"],
            ["Simulation console UI", "Roadmap", "API-first today; UI next sprint"],
            ["Postgres ledger", "Roadmap", "Budgeted in pilot"],
            ["Golden Edge v3 branding", "Policy framework name", "Map rules to JSON schema"],
        ],
        [44, 36, 94],
    )
    pdf.section("Funding & Commercial Readiness")
    pdf.bullets([
        "IRAP-ready narrative: complete (see grant-core-narrative.pdf)",
        "LOI template: ready for first pilot partner",
        "Budget model: $120K pilot / $60K funding ask",
        "Commercial path: $10K -> $50K -> $120K",
        "Website alignment: noetfield.com messaging matches GEL positioning",
    ])
    pdf.section("Next 30 Days - Internal Priority")
    pdf.bullets([
        "Deploy api.noetfield.com with API key auth",
        "Postgres migration + audit retention policy",
        "First LOI signed + sandbox pilot kickoff",
        "Trust Ledger export from audit rows",
        "Submit IRAP expression of interest",
    ])


# ---------------------------------------------------------------------------
# Slide deck builder
# ---------------------------------------------------------------------------

def _slide_title(slide, title: str, subtitle: str = "") -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*NOETFIELD_DARK)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(8.8), Inches(1))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(180, 200, 230)


def _slide_content(slide, title: str, bullets: list[str]) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.8))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(*NOETFIELD_BLUE)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(35, 40, 48)
        p.space_after = Pt(12)


def build_pitch_deck(audience: str) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides_data: list[tuple[str, list[str] | str]] = [
        ("Noetfield Governance Execution Layer", "Pre-execution governance for regulated AI systems"),
        ("The Problem", [
            "AI adoption is outpacing governance capacity",
            "No pre-execution validation before operational actions",
            "No audit-ready traceability for board oversight",
        ]),
        ("The Risk", [
            "Regulatory exposure under emerging AI rules",
            "Board liability from invisible decision chains",
            "Operational blind spots in automated workflows",
        ]),
        ("The Solution", [
            "Governance Execution Layer (GEL)",
            "Evaluate intent BEFORE external systems execute",
            "Non-custodial: governance signals only",
        ]),
        ("The Product - Golden Edge v3", [
            "ALLOW / REJECT / REVIEW decision engine",
            "Policy enforcement + corridor rules",
            "Immutable audit ledger",
        ]),
        ("Console Demo Flow", [
            "Intent input (applicant / AI decision payload)",
            "Governed decision output with score breakdown",
            "Audit trail persisted for compliance review",
        ]),
        ("Architecture", [
            "FastAPI governance API",
            "Postgres-backed event ledger (pilot target)",
            "Policy loader + Trust Ledger export path",
        ]),
        ("Pilot Scope", [
            "6-week structured engagement",
            "BC organization - read-only integration",
            "Simulation + policy mapping + audit review",
        ]),
        ("Why Canada", [
            "AI governance leadership opportunity",
            "Regulatory alignment (AI Act readiness)",
            "Safe adoption layer for banking & public sector",
        ]),
        ("The Ask", [
            "IRAP / BC Innovation co-funding ($60K match)",
            "Pilot partner for enterprise validation",
            "Board-level governance modernization",
        ]),
    ]

    if audience == "INTERNAL":
        slides_data.append(("Internal - Reality & Commercial Path", [
            "Built: FastAPI engine, policy JSON, SQLite audit, portal API",
            "Next: Postgres, console UI, Trust Ledger export, api.noetfield.com",
            "Commercial ladder: $10K entry -> $50K pilot -> $120K IRAP program",
        ]))

    for i, item in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        if i == 0:
            _slide_title(slide, item[0], item[1])  # type: ignore[arg-type]
        else:
            _slide_content(slide, item[0], item[1])  # type: ignore[arg-type]

    out_dir = OUT / audience.lower()
    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / "pitch-deck-noetfield-gel.pptx"
    prs.save(str(path))
    return path


def build_pitch_deck_pdf(audience: str) -> Path:
    """Handout version of the deck (one page per slide summary)."""
    slides = [
        ("Slide 1 - Title", "Noetfield Governance Execution Layer"),
        ("Slide 2 - Problem", "AI adoption without governance; no pre-execution validation; no audit traceability."),
        ("Slide 3 - Risk", "Regulatory exposure; board liability; invisible AI decision chains."),
        ("Slide 4 - Solution", "Governance Execution Layer - evaluate before execution."),
        ("Slide 5 - Product", "Golden Edge v3: ALLOW/REJECT engine, policy enforcement, immutable ledger."),
        ("Slide 6 - Console Demo", "Intent input -> decision output -> audit trail."),
        ("Slide 7 - Architecture", "FastAPI + Postgres event ledger + policy engine."),
        ("Slide 8 - Pilot Scope", "6 weeks; BC org; read-only integration."),
        ("Slide 9 - Why Canada", "AI governance leadership; regulatory alignment; safe adoption."),
        ("Slide 10 - Ask", "IRAP funding; pilot partner; enterprise validation."),
    ]
    if audience == "INTERNAL":
        slides.append(("Slide 11 - Internal", "Build status, roadmap, commercial ladder $10K/$50K/$120K."))

    meta = DocMeta(
        filename="pitch-deck-handout.pdf",
        title="Pitch Deck Handout",
        subtitle="Noetfield Governance Execution Layer - slide summaries",
        audience=audience,
        classification="Presentation companion",
    )

    def builder(pdf: NoetfieldPDF) -> None:
        pdf.cover()
        for title, body in slides:
            pdf.section(title)
            pdf.body(body)

    return save_pdf(meta, builder)


def main() -> None:
    pdfs = [
        (DocMeta("grant-core-narrative.pdf", "Grant Core Narrative", "IRAP-Grade Positioning - Governance Execution Layer", "EXTERNAL", "Grant submission"), build_grant_narrative_external),
        (DocMeta("grant-core-narrative.pdf", "Grant Core Narrative", "IRAP-Grade Positioning + Internal Build Appendix", "INTERNAL", "Grant + engineering"), build_grant_narrative_internal),
        (DocMeta("loi-template.pdf", "Letter of Intent", "Governance Execution Pilot - Partner LOI", "EXTERNAL", "Partner template"), build_loi_external),
        (DocMeta("loi-template.pdf", "Letter of Intent", "Governance Execution Pilot - LOI + Sales Notes", "INTERNAL", "Sales enablement"), build_loi_internal),
        (DocMeta("budget-breakdown.pdf", "Budget Breakdown", "IRAP-Style Project Budget - 6-Week Pilot", "EXTERNAL", "Funding request"), build_budget_external),
        (DocMeta("budget-breakdown.pdf", "Budget Breakdown", "IRAP Budget + Commercial Ladder", "INTERNAL", "Finance + sales"), build_budget_internal),
        (DocMeta("strategic-positioning.pdf", "Strategic Positioning", "Category Definition - Governance Execution Infrastructure", "EXTERNAL", "Market positioning"), build_strategic_external),
        (DocMeta("strategic-positioning.pdf", "Strategic Positioning", "Category + Moat + Messaging Guardrails", "INTERNAL", "Strategy"), build_strategic_internal),
        (DocMeta("reality-check.pdf", "Reality Check", "Current Build Status & 30-Day Priorities", "INTERNAL", "Confidential - ops"), build_reality_check_internal),
    ]

    generated: list[Path] = []
    for meta, builder in pdfs:
        path = save_pdf(meta, builder)
        generated.append(path)
        print(f"PDF  -> {path}")

    for aud in ("EXTERNAL", "INTERNAL"):
        pptx = build_pitch_deck(aud)
        pdf = build_pitch_deck_pdf(aud)
        generated.append(pptx)
        generated.append(pdf)
        print(f"PPTX -> {pptx}")
        print(f"PDF  -> {pdf}")

    print(f"\nDone. {len(generated)} files in docs/output/")


if __name__ == "__main__":
    main()
